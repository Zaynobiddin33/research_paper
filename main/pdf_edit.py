from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import qrcode
from io import BytesIO
import re
import os
from datetime import datetime


def give_certificate(first_name, last_name, article_name, url):
    username = f"{first_name} {last_name}"
    data = {
        "user_name": username,
        'article_name': article_name,
        'num_years' : str(datetime.now().year - 2024),
        'num_month' : str(datetime.now().month),
        'publish year': str(datetime.now().year),
        'today' :datetime.now().strftime("%d.%m.%Y"),
        'paper_url': url
        
    }
    prs = Presentation("main/tmplt.pptx")

    # === Generate QR code ===
    qr = qrcode.QRCode(box_size=10, border=2)
    qr.add_data(data["paper_url"])
    qr.make(fit=True)
    qr_img = qr.make_image(fill_color="black", back_color="white")

    qr_buffer = BytesIO()
    qr_img.save(qr_buffer, format="PNG")
    qr_buffer.seek(0)

    # === Regex ===
    PLACEHOLDER_PATTERN = re.compile(r"\{\{(.+?)\}\}")

    # === Helpers ===
    def get_run_style(run):
        font = run.font
        style = {
            "name": font.name,
            "size": font.size,
            "bold": font.bold,
            "italic": font.italic,
            "underline": font.underline,
            "color": None
        }
        try:
            if font.color and font.color.rgb:
                style["color"] = font.color.rgb
        except Exception:
            pass
        return style


    def apply_style_to_run(run, style):
        font = run.font
        if style["name"]:
            font.name = style["name"]
        if style["size"]:
            font.size = style["size"]
        if style["bold"] is not None:
            font.bold = style["bold"]
        if style["italic"] is not None:
            font.italic = style["italic"]
        if style["underline"] is not None:
            font.underline = style["underline"]
        if style["color"]:
            try:
                font.color.rgb = RGBColor(style["color"][0], style["color"][1], style["color"][2])
            except Exception:
                pass


    def replace_placeholders_in_paragraph(paragraph, replacements, slide, shape):
        """Replace placeholders and insert QR code at text location."""
        combined_text = "".join(r.text for r in paragraph.runs)
        if "{{" not in combined_text:
            return

        # Handle QR placeholder
        if "{{QR_CODE}}" in combined_text:
            # Clear text
            paragraph.clear()

            # === Place QR exactly at shape's position ===
            left = shape.left
            top = shape.top
            width = Inches(1.5)   # smaller size
            height = Inches(1.5)

            # Remove the original shape (the text box with placeholder)
            slide.shapes._spTree.remove(shape._element)

            # Add QR picture at same spot
            qr_buffer.seek(0)
            slide.shapes.add_picture(qr_buffer, left, top, width=width, height=height)
            print(f"✅ Inserted QR at position ({left}, {top})")
            return

        # Replace text placeholders normally
        def replace_fn(m):
            key = m.group(1)
            return replacements.get(key, m.group(0))
        new_text = PLACEHOLDER_PATTERN.sub(replace_fn, combined_text)

        if new_text == combined_text:
            return

        styles = [get_run_style(r) for r in paragraph.runs]
        paragraph.clear()
        new_run = paragraph.add_run()
        new_run.text = new_text
        if styles:
            apply_style_to_run(new_run, styles[0])


    # === Replace placeholders across slides ===
    for slide in prs.slides:
        # Copy the list of shapes because we might remove some while iterating
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                replace_placeholders_in_paragraph(paragraph, data, slide, shape)

    # === Save final result ===
    if not os.path.exists("media/certificates"):
        os.makedirs('media/certificates')
    name = article_name.replace(' ', "-")
    prs.save(f"media/certificates/{first_name}-{last_name}-{name}.pptx")
    print("file created")
    
    return f"certificates/{first_name}-{last_name}-{name}.pptx"